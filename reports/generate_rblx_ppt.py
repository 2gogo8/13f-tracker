#!/usr/bin/env python3
"""Generate single-company deep research PPT for RBLX (Roblox Corporation)
   v3 Style: Black background, RED section titles, gold accents, pattern score badge
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors (v3 style)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
CARD_BG = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xF0, 0xF0, 0xF0)
GRAY = RGBColor(0x99, 0x99, 0x99)
DIM = RGBColor(0x66, 0x66, 0x66)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
RED = RGBColor(0xC4, 0x1E, 0x3A)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'PingFang TC'
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.width = Pt(0.5)
    return shape

def add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def watermark(slide):
    add_text(slide, 10.5, 7.0, 2.5, 0.3, 'JG的反市場報告書', 8, DIM, align=PP_ALIGN.RIGHT)

def slide_num(slide, num):
    add_text(slide, 0.3, 7.0, 0.5, 0.3, str(num), 8, DIM)

# ============================================================
# SLIDE 1: COVER
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

# Gradient effect background
gradient_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(0.5), Inches(6.5), Inches(3.5))
gradient_circle.fill.solid()
gradient_circle.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
gradient_circle.fill.transparency = 0.6
gradient_circle.line.fill.background()

add_text(s, 1.5, 1.5, 10, 1.2, "ROBLOX CORPORATION", 54, GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 2.8, 10, 0.6, "深度研究報告 — 元宇宙平台的盈利拐點", 20, GRAY, align=PP_ALIGN.CENTER)
add_gold_line(s, 5.5, 3.6, 2.3)

# Pattern score badge with glow
x = 5.5
# Glow effect
glow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.05), Inches(4.1), Inches(2.5), Inches(1.3))
glow.fill.solid()
glow.fill.fore_color.rgb = GOLD
glow.fill.transparency = 0.85
glow.line.fill.background()

# Main badge card
card = add_card(s, x, 4.15, 2.4, 1.2)
add_text(s, x, 4.2, 2.4, 0.3, 'RBLX', 16, GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, x, 4.52, 2.4, 0.5, 'C', 36, GRAY, bold=True, align=PP_ALIGN.CENTER)  # Grade C
add_text(s, x, 4.95, 2.4, 0.25, '48', 14, WHITE, align=PP_ALIGN.CENTER)  # Score 48
add_text(s, x, 5.2, 2.4, 0.15, '型態評級', 10, DIM, align=PP_ALIGN.CENTER)

add_text(s, 1.5, 5.7, 10, 0.4, "2026年2月19日", 14, DIM, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 6.2, 10, 0.3, "僅供參考，不構成投資建議", 10, DIM, align=PP_ALIGN.CENTER)
watermark(s)
slide_num(s, 1)

# ============================================================
# SLIDE 2: TITLE PAGE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

# Decorative gold dots
for dot_x in [0.5, 12.5]:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dot_x), Inches(1.5), Pt(8), Pt(8))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.fill.transparency = 0.3
    dot.line.fill.background()

add_text(s, 0.8, 1.0, 10, 1.0, "RBLX", 64, GOLD, bold=True)
add_text(s, 0.8, 2.2, 10, 0.5, "Roblox Corporation", 20, GRAY)
add_gold_line(s, 0.8, 2.9, 2)
add_text(s, 0.8, 3.2, 10, 0.4, "$63.05 · 市值 $447億美元", 18, WHITE)

# Key stats boxes
stats = [
    ('2025 營收', '$48.9B', GOLD),
    ('YoY 成長', '+35.8%', GREEN),
    ('淨虧損', '-$10.7B', RED),
    ('DAU', '88.9M', BLUE),
]
for i, (label, value, color) in enumerate(stats):
    x = 0.8 + i * 2.8
    add_card(s, x, 4.0, 2.5, 1.0)
    add_text(s, x + 0.15, 4.05, 2.2, 0.55, value, 28, color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 4.6, 2.2, 0.3, label, 10, GRAY, align=PP_ALIGN.CENTER)

watermark(s)
slide_num(s, 2)

# ============================================================
# SLIDE 3: COMPANY OVERVIEW
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 5, 0.6, "RBLX 公司概況", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

overview_lines = [
    '▸ 全球最大 UGC (用戶生成內容) 遊戲平台',
    '▸ CEO: David Baszucki · 員工: 2,474人 · 2021年3月上市',
    '',
    '▸ 核心產品:',
    '    Roblox Platform — 人類共同體驗平台',
    '    Roblox Studio — 免費創作工具，1,500萬創作者',
    '    Roblox Client — 應用程式，4,000萬個遊戲體驗',
    '',
    '▸ 2025年亮點:',
    '    營收 $48.9億 (+35.8%)，成長重新加速',
    '    DAU 8,890萬 (+11.8%)，用戶黏性持續強化',
    '    營業虧損率從 -45% (2023) 收窄至 -25.2%',
    '    自由現金流 $13.5億 (+111%)，朝向盈虧平衡穩步前進',
]
for i, line in enumerate(overview_lines):
    c = WHITE if not line.startswith('▸') else RGBColor(0xCC, 0xCC, 0xCC)
    add_text(s, 0.8, 1.3 + i * 0.35, 11.5, 0.35, line, 13, c)

watermark(s)
slide_num(s, 3)

# ============================================================
# SLIDE 4: BUSINESS MODEL
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 商業模式 — Bookings vs Revenue", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

model_lines = [
    '## Roblox 的財務架構有別於傳統遊戲公司',
    '',
    '★ Bookings（預訂）:',
    '▸ 用戶購買虛擬貨幣「Robux」的現金流入',
    '▸ 包含：直接購買 + Premium 訂閱費用',
    '▸ 這是衡量平台健康度的核心指標（類似現金收入）',
    '',
    '★ Revenue（營收）:',
    '▸ 當 Robux 被消費時才認列收入',
    '▸ 由於用戶囤積 Robux，收入認列存在時間差',
    '▸ 遞延收入反映在資產負債表',
    '',
    '★ 2025年數據對比:',
    '▸ Bookings: ~$54-56億美元（估計）',
    '▸ Revenue: $48.9億美元',
    '▸ 差異原因: 用戶囤積 + 平台成長期遞延收入累積',
    '',
    '+投資者應關注: Bookings 成長率 > Revenue 成長率',
    '  代表用戶消費意願強勁，未來收入能見度高',
]
for i, line in enumerate(model_lines):
    c = GOLD if line.startswith('##') or line.startswith('★') else (GREEN if line.startswith('+') else WHITE)
    fs = 15 if line.startswith('##') else 13
    txt = line.lstrip('##★+▸ ')
    add_text(s, 0.8, 1.3 + i * 0.30, 11.5, 0.30, txt, fs, c, bold=line.startswith('★') or line.startswith('##'))

watermark(s)
slide_num(s, 4)

# ============================================================
# SLIDE 5: USER METRICS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 用戶數據與成長趨勢", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

metrics = [
    ('平均 DAU', '88.9M', '+11.8%', GREEN),
    ('總遊玩時數', '812億小時', '+14.2%', BLUE),
    ('ARPDAU', '$16.90', '+7.0%', YELLOW),
    ('13歲以上占比', '57%', '成熟化', WHITE),
]
for i, (label, value, change, color) in enumerate(metrics):
    y = 1.4 + i * 1.2
    add_card(s, 0.8, y, 11.5, 1.0)
    add_text(s, 1.0, y + 0.05, 3, 0.4, label, 14, GRAY, bold=True)
    add_text(s, 1.0, y + 0.5, 3, 0.4, value, 24, color, bold=True)
    add_text(s, 5.5, y + 0.3, 6, 0.4, change, 18, color, bold=True)

watermark(s)
slide_num(s, 5)

# ============================================================
# SLIDE 6: REVENUE STRUCTURE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 5, 0.6, "RBLX 營收結構", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

rev_items = [
    ('核心遊戲內消費', 92, '~$45B (92%)', GOLD),
    ('Premium 訂閱', 7, '~$3.4B (7%)', BLUE),
    ('廣告收入 (新興)', 1, '~$50-80M (1%)', GREEN),
]
for i, (label, pct, amount, color) in enumerate(rev_items):
    y = 1.6 + i * 0.7
    add_text(s, 0.8, y, 2.5, 0.3, label, 13, GRAY, align=PP_ALIGN.RIGHT)
    bar_w = pct / 100 * 8
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(y), Inches(max(0.5, bar_w)), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(s, 3.5 + bar_w + 0.2, y + 0.05, 3, 0.3, amount, 11, WHITE)

add_text(s, 0.8, 4.0, 11.5, 0.35, '★ 2026年展望: 廣告收入可能快速成長至 $200-300M（占比3-5%）', 14, GOLD, bold=True)
add_text(s, 0.8, 4.4, 11.5, 0.35, '▸ 電商與品牌合作可能達到 $100-150M', 13, WHITE)

watermark(s)
slide_num(s, 6)

# ============================================================
# SLIDE 7: PATH TO PROFITABILITY (3 SCENARIOS)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 獲利路徑分析 — 三種情境", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

scenarios = [
    {
        'title': '保守路徑 (2028年盈虧平衡)',
        'color': YELLOW,
        'points': [
            '營收年均成長 25%',
            'R&D/營收 32%→28%',
            '2028年營業利潤率 +3%',
        ]
    },
    {
        'title': '中性路徑 (2027年盈虧平衡) ★ 最可能',
        'color': GREEN,
        'points': [
            '營收年均成長 30%',
            '廣告收入 2027年達 $5億',
            '2027年營業利潤率 +2%',
        ]
    },
    {
        'title': '樂觀路徑 (2026年盈虧平衡)',
        'color': GOLD,
        'points': [
            '營收年均成長 35%',
            'AI生成內容降低開發成本',
            '2026年營業利潤率 +0.8%',
        ]
    },
]

for i, sc in enumerate(scenarios):
    y = 1.5 + i * 1.6
    add_card(s, 0.8, y, 11.5, 1.4)
    add_text(s, 1.0, y + 0.05, 10, 0.35, sc['title'], 15, sc['color'], bold=True)
    for j, pt in enumerate(sc['points']):
        add_text(s, 1.0, y + 0.45 + j * 0.3, 10, 0.3, f"▸ {pt}", 12, WHITE)

watermark(s)
slide_num(s, 7)

# ============================================================
# SLIDE 8: CAPEX APPLICATION SCENARIOS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 資本支出應用場景（白話解釋）", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

capex_lines = [
    '★ 2025年 CapEx $4.41億美元（占營收9%）',
    '',
    '錢花在哪裡？',
    '',
    '▸ 資料中心與雲端基礎設施 (60%, $2.6億):',
    '  → 8,890萬用戶同時在線，租用AWS、Google Cloud',
    '  → 就像開全球連鎖餐廳，需要在各地蓋廚房與物流中心',
    '',
    '▸ AI運算資源 (20%, $0.9億):',
    '  → 訓練AI助手、內容審核AI（每天審查數億條訊息）',
    '  → 聘請24小時不休息的AI保全團隊',
    '',
    '▸ 辦公室與設備 (15%, $0.7億):',
    '  → 員工從2,100人增至2,474人，擴張辦公空間',
    '',
    '▸ 其他技術設施 (5%, $0.2億):',
    '  → 網路安全設備、測試裝置（手機、VR頭盔）',
    '',
    '!CapEx/營收比 9% 屬於成長型科技公司正常水準',
    '+預期2027-2028年降至6-7%',
]
for i, line in enumerate(capex_lines):
    c = GOLD if line.startswith('★') else (RED if line.startswith('!') else (GREEN if line.startswith('+') else (BLUE if line.startswith('▸') else WHITE)))
    txt = line.lstrip('★!+▸')
    add_text(s, 0.8, 1.3 + i * 0.3, 11.5, 0.3, txt, 12, c, bold=line.startswith('★'))

watermark(s)
slide_num(s, 8)

# ============================================================
# SLIDE 9: 2026 MAJOR MOVES
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 2026年重大動向", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

moves = [
    ('+ 廣告平台全面推出 (Q2)', '向所有開發者開放廣告API，目標$2-3億收入', GREEN),
    ('+ AI創作工具大升級', 'Roblox Assistant自然語言生成3D物件，降低創作門檻', GREEN),
    ('+ 電商化戰略', '虛擬商品實體化，與Nike、Gucci等品牌合作', GREEN),
    ('+ 17+內容拓展', '開放成人內容，吸引18-25歲高消費用戶群', GREEN),
    ('', '', WHITE),
    ('- 歐盟DSA法規', '可能增加合規成本$50-100M/年', RED),
    ('- 中國市場停滯', '2022年下架，短期無重啟跡象', RED),
    ('- 兒童保護爭議', '平台安全與監管風險持續', RED),
]
for i, (title, desc, color) in enumerate(moves):
    if title:
        y = 1.4 + i * 0.62
        add_card(s, 0.8, y, 11.5, 0.55)
        add_text(s, 1.0, y + 0.05, 10, 0.25, title, 13, color, bold=True)
        add_text(s, 1.0, y + 0.3, 10.5, 0.2, desc, 11, WHITE)

watermark(s)
slide_num(s, 9)

# ============================================================
# SLIDE 10: CATALYST TIMELINE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 2026年催化劑時間表", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

catalysts = [
    ('2月', 'Q4 2025財報公布', '已超預期，營收+36%'),
    ('5月', 'Q1 2026財報', '關鍵:DAU成長、廣告收入首次披露'),
    ('6月', '廣告平台正式推出', '最大催化劑，對標Unity Ads'),
    ('8月', 'Q2 2026財報', '廣告收入規模、17+內容影響'),
    ('9月', 'Roblox Developers Conference', '新工具發布、2027戰略指引'),
    ('12月', '年終購物季', 'Q4佔全年Bookings 30-35%'),
]
for i, (time, event, impact) in enumerate(catalysts):
    y = 1.5 + i * 0.7
    add_card(s, 0.8, y, 11.5, 0.6)
    add_text(s, 1.0, y + 0.05, 1.5, 0.4, time, 13, GOLD, bold=True)
    add_text(s, 2.8, y + 0.05, 5, 0.4, event, 14, WHITE, bold=True)
    c = GREEN if '催化' in impact or '超預期' in impact else GRAY
    add_text(s, 8.5, y + 0.05, 3, 0.4, impact, 11, c, align=PP_ALIGN.RIGHT)

watermark(s)
slide_num(s, 10)

# ============================================================
# SLIDE 11: COMPETITOR COMPARISON
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 競爭對手比較", 30, BLUE, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

competitors = [
    {
        'name': 'RBLX vs Fortnite (Epic Games)',
        'strategy': 'Fortnite: 高ARPDAU ($25)，單一遊戲依賴 | RBLX: 1,500萬創作者生態，ARPDAU偏低 ($16.90)',
        'advantage': '✅ RBLX優勢: 創作者生態150倍於Fortnite ❌ 劣勢: ARPDAU落後33%',
    },
    {
        'name': 'RBLX vs Minecraft (Microsoft)',
        'strategy': 'Minecraft: 買斷制穩定收入，教育市場強 | RBLX: 免費+虛擬商品，社交功能強',
        'advantage': '✅ RBLX優勢: 社交平台屬性，用戶黏性更高（2.5小時/天）',
    },
    {
        'name': 'RBLX vs Meta Horizon Worlds',
        'strategy': 'Meta: VR社交平台，硬體依賴 | RBLX: 跨平台，手機/PC/主機全覆蓋',
        'advantage': '✅ RBLX優勢: 用戶基數89M vs Meta <1M，跨平台優勢明顯',
    },
]

for i, comp in enumerate(competitors):
    y = 1.5 + i * 1.5
    add_card(s, 0.8, y, 11.5, 1.35)
    add_text(s, 1.0, y + 0.05, 10, 0.35, comp['name'], 15, GOLD, bold=True)
    add_text(s, 1.0, y + 0.45, 10.5, 0.4, comp['strategy'], 12, WHITE)
    add_text(s, 1.0, y + 0.9, 10.5, 0.35, comp['advantage'], 11, GREEN)

watermark(s)
slide_num(s, 11)

# ============================================================
# SLIDE 12: INVESTMENT VERDICT (最大優勢 + 最大難處)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

# Decorative gold accent lines
for line_x in [0.5, 12.7]:
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(line_x), Inches(1.4), Pt(2), Inches(4.5))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.fill.transparency = 0.5
    line.line.fill.background()

add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 投資結論", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

# Left: 最大優勢
add_text(s, 0.8, 1.4, 5.5, 0.4, "最大優勢", 16, GREEN, bold=True)
advantages = [
    '全球最大UGC平台，1,500萬創作者',
    '用戶黏性極高，每日2.5小時遊玩',
    '營收成長重新加速至36% (2025)',
    '虧損持續收窄，2027年可能盈虧平衡',
    '廣告+電商新收入源即將爆發',
]
for i, adv in enumerate(advantages):
    add_text(s, 0.8, 1.9 + i * 0.4, 5.5, 0.4, f"▸ {adv}", 12, RGBColor(0xCC, 0xCC, 0xCC))

# Right: 最大難處
add_text(s, 6.8, 1.4, 5.5, 0.4, "最大難處", 16, RED, bold=True)
challenges = [
    '仍未盈利，2025虧損$10.7億',
    'ARPDAU僅$16.90，落後Fortnite 33%',
    'DAU成長放緩至12%，天花板隱現',
    '兒童保護爭議，監管風險高',
    '估值P/S 9.1x偏高，回調風險大',
]
for i, ch in enumerate(challenges):
    add_text(s, 6.8, 1.9 + i * 0.4, 5.5, 0.4, f"▸ {ch}", 12, RGBColor(0xCC, 0xCC, 0xCC))

# Bottom: Key message
add_card(s, 0.8, 5.8, 11.5, 0.9)
key_msg = '2027年盈虧平衡是關鍵拐點 — 廣告業務若不成功，ARPDAU難有突破\n適合長期持有（3-5年），但需承受高波動風險'
add_text(s, 1.0, 5.85, 11, 0.8, key_msg, 14, GOLD, align=PP_ALIGN.CENTER)

watermark(s)
slide_num(s, 12)

# ============================================================
# SLIDE 13: MONITORING KPIs
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "RBLX 關鍵監控指標（KPIs）", 30, RED, bold=True)
add_gold_line(s, 0.8, 1.0, 1.5)

kpis = [
    ('成長指標', [
        'DAU YoY成長率 (目標: >10%)',
        'Hours Engaged成長率 (目標: >12%)',
        'ARPDAU YoY成長率 (目標: >8%)',
    ], GOLD),
    ('獲利能力', [
        '營業利潤率 (目標: 2026年-15%、2027年0%)',
        '自由現金流 (目標: 維持正值)',
        'R&D費用/營收比 (目標: 降至28%以下)',
    ], GREEN),
    ('新業務', [
        '廣告收入規模 (目標: 2026年$2-3億)',
        '17+內容佔比 (目標: 2026年達10%)',
        'AI創作工具採用率',
    ], BLUE),
    ('風險信號', [
        'DAU成長跌破8% → 減倉信號',
        '營業虧損率擴大 → 立即停損',
        '重大監管訴訟 → 評估退出',
    ], RED),
]

y_start = 1.4
for cat, items, color in kpis:
    add_text(s, 0.8, y_start, 11.5, 0.4, cat, 15, color, bold=True)
    for i, item in enumerate(items):
        add_text(s, 1.2, y_start + 0.4 + i * 0.3, 11, 0.3, f"□ {item}", 12, WHITE)
    y_start += 0.4 + len(items) * 0.3 + 0.3

watermark(s)
slide_num(s, 13)

# ============================================================
# SAVE
# ============================================================
output_path = '/Users/jgtruestock/.openclaw/workspace/projects/13f-tracker/reports/RBLX_深度研究.pptx'
prs.save(output_path)
print(f"✅ PPT generated: {output_path}")
print(f"📄 Total slides: {len(prs.slides)}")

import os
file_size = os.path.getsize(output_path)
print(f"📦 File size: {file_size / 1024:.1f} KB")
