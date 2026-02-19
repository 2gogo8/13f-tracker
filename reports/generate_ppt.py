#!/usr/bin/env python3
"""Generate 50+ page black-background PPT for JG's Anti-Market Report - Version 2.0
   Changes:
   1. CapEx slides → "資本支出應用場景" with real-world analogies
   2. Verdict slide → "投資結論" with advantages/challenges (no target price)
   3. New slide: "競爭對手比較" for each company
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DARK_BG = RGBColor(0x12, 0x12, 0x12)
CARD_BG = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xF0, 0xF0, 0xF0)
GRAY = RGBColor(0x99, 0x99, 0x99)
DIM = RGBColor(0x66, 0x66, 0x66)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
RED = RGBColor(0xC4, 0x1E, 0x3A)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='PingFang TC'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_para(tf, text, font_size=14, color=WHITE, bold=False, space_before=0, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'PingFang TC'
    p.space_before = Pt(space_before)
    p.alignment = align
    return p

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.width = Pt(0.5)
    return shape

def add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_bar(slide, left, top, width, height, color, max_width=5.0, value=1.0, max_val=1.0):
    bar_w = max(0.3, width * (value / max_val))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(bar_w), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def watermark(slide):
    add_text(slide, 10.5, 7.0, 2.5, 0.3, 'JG的反市場報告書', 8, DIM, align=PP_ALIGN.RIGHT)

def slide_num(slide, num):
    add_text(slide, 0.3, 7.0, 0.5, 0.3, str(num), 8, DIM)

# ============================================================
# SLIDE 1: COVER
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(s)

# Subtle radial gradient effect - large semi-transparent dark circle behind title for depth
gradient_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(0.5), Inches(6.5), Inches(3.5))
gradient_circle.fill.solid()
gradient_circle.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
gradient_circle.fill.transparency = 0.6
gradient_circle.line.fill.background()

add_text(s, 1.5, 1.5, 10, 1.2, "JG的反市場報告書", 54, GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 2.8, 10, 0.6, "深度研究系列 — 五檔反市場精選個股", 20, GRAY, align=PP_ALIGN.CENTER)
add_gold_line(s, 5.5, 3.6, 2.3)

# 5 badges - Pattern Score
pattern_data = [
    ('GOOGL', 'B', 62, GOLD),
    ('PLTR', 'A', 80, GREEN),
    ('RKLB', 'C', 48, GRAY),
    ('AMD', 'B', 58, GOLD),
    ('ARES', 'B', 67, GOLD),
]
for i, (ticker, grade, score, grade_color) in enumerate(pattern_data):
    x = 1.5 + i * 2.2
    # Glow effect - semi-transparent gold rounded rectangle behind badge
    glow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.05), Inches(4.1), Inches(2.3), Inches(1.3))
    glow.fill.solid()
    glow.fill.fore_color.rgb = GOLD
    glow.fill.transparency = 0.85
    glow.line.fill.background()
    
    # Main badge card
    card = add_card(s, x, 4.15, 2.2, 1.2)
    
    # Ticker at top
    add_text(s, x, 4.2, 2.2, 0.3, ticker, 16, GOLD, bold=True, align=PP_ALIGN.CENTER)
    
    # Grade letter (large, color-coded)
    add_text(s, x, 4.52, 2.2, 0.5, grade, 36, grade_color, bold=True, align=PP_ALIGN.CENTER)
    
    # Score number
    add_text(s, x, 4.95, 2.2, 0.25, str(score), 14, WHITE, align=PP_ALIGN.CENTER)
    
    # Label at bottom
    add_text(s, x, 5.2, 2.2, 0.15, '型態評級', 10, DIM, align=PP_ALIGN.CENTER)

add_text(s, 1.5, 5.5, 10, 0.4, "2026年2月19日", 14, DIM, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 6.0, 10, 0.3, "僅供參考，不構成投資建議", 10, DIM, align=PP_ALIGN.CENTER)
watermark(s)
slide_num(s, 1)

# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 5, 0.7, "目錄", 36, GOLD, bold=True)
add_gold_line(s, 0.8, 1.2, 1.5)

toc = [
    ("01", "GOOGL — Alphabet Inc.", "成長瓶頸分析 + 2026突破點", "P.3-11"),
    ("02", "PLTR — Palantir Technologies", "成長持續性 + 利潤率 + 2027方向", "P.12-20"),
    ("03", "RKLB — Rocket Lab USA", "未來方向 + 政策面 + 成長率", "P.21-29"),
    ("04", "AMD — Advanced Micro Devices", "低利潤率突破路徑", "P.30-38"),
    ("05", "ARES — Ares Management", "政策影響 + 寬鬆/緊縮情境", "P.39-47"),
    ("06", "總結與比較", "五檔精選綜合評估", "P.48-49"),
]
for i, (num, title, sub, page) in enumerate(toc):
    y = 1.8 + i * 0.85
    add_text(s, 0.8, y, 0.6, 0.5, num, 24, GOLD, bold=True)
    add_text(s, 1.6, y, 6, 0.35, title, 18, WHITE, bold=True)
    add_text(s, 1.6, y + 0.35, 6, 0.3, sub, 12, GRAY)
    add_text(s, 10.5, y + 0.1, 2, 0.3, page, 12, DIM, align=PP_ALIGN.RIGHT)

watermark(s)
slide_num(s, 2)

# ============================================================
# Helper: company section generator
# ============================================================
page_counter = [3]

def make_company_section(ticker, name, subtitle, price_str, sigma, data):
    """Generate 9-10 slides per company"""
    pc = page_counter[0]
    
    # --- SLIDE: Title ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    
    # Decorative gold dots as accent shapes
    for dot_x in [0.5, 12.5]:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dot_x), Inches(1.5), Pt(8), Pt(8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.fill.transparency = 0.3
        dot.line.fill.background()
    
    add_text(s, 0.8, 1.0, 10, 1.0, ticker, 64, GOLD, bold=True)
    add_text(s, 0.8, 2.2, 10, 0.5, name, 20, GRAY)
    add_gold_line(s, 0.8, 2.9, 2)
    add_text(s, 0.8, 3.2, 10, 0.4, f"{price_str}  ·  σ{sigma}", 18, WHITE)
    # Key stats boxes
    stats = data.get('key_stats', [])
    for i, (label, value, color) in enumerate(stats):
        x = 0.8 + i * 2.8
        add_card(s, x, 4.0, 2.5, 1.0)
        add_text(s, x + 0.15, 4.05, 2.2, 0.55, value, 28, color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.15, 4.6, 2.2, 0.3, label, 10, GRAY, align=PP_ALIGN.CENTER)
    watermark(s); slide_num(s, pc); pc += 1
    
    # --- SLIDE: Company Overview ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, 0.8, 0.4, 5, 0.6, f"{ticker} 公司概況", 30, RED, bold=True)
    add_gold_line(s, 0.8, 1.0, 1.5)
    overview = data.get('overview', [])
    for i, line in enumerate(overview):
        add_text(s, 0.8, 1.3 + i * 0.35, 11, 0.35, line, 13, WHITE if not line.startswith('▸') else RGBColor(0xCC, 0xCC, 0xCC))
    watermark(s); slide_num(s, pc); pc += 1

    # --- SLIDE: Revenue Breakdown ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, 0.8, 0.4, 5, 0.6, f"{ticker} 營收結構", 30, RED, bold=True)
    add_gold_line(s, 0.8, 1.0, 1.5)
    rev_items = data.get('revenue', [])
    for i, (label, pct, amount, color) in enumerate(rev_items):
        y = 1.4 + i * 0.65
        add_text(s, 0.8, y, 2, 0.3, label, 13, GRAY, align=PP_ALIGN.RIGHT)
        bar_w = pct / 100 * 8
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(y), Inches(max(0.5, bar_w)), Inches(0.4))
        shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
        add_text(s, 3.0 + bar_w + 0.2, y + 0.05, 3, 0.3, amount, 11, WHITE)
    watermark(s); slide_num(s, pc); pc += 1
    
    # --- SLIDE: CapEx 資本支出應用場景 (CHANGED) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, 0.8, 0.4, 8, 0.6, f"{ticker} 資本支出應用場景", 30, RED, bold=True)
    add_gold_line(s, 0.8, 1.0, 1.5)
    capex = data.get('capex_scenarios', [])
    for i, line in enumerate(capex):
        c = GOLD if line.startswith('★') else (RED if line.startswith('!') else (BLUE if line.startswith('▸') else WHITE))
        add_text(s, 0.8, 1.3 + i * 0.35, 11.5, 0.35, line.lstrip('★!▸'), 13, c)
    watermark(s); slide_num(s, pc); pc += 1

    # --- SLIDE: 2026 Major Moves ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, 0.8, 0.4, 8, 0.6, f"{ticker} 2026 重大動向與消息", 30, RED, bold=True)
    add_gold_line(s, 0.8, 1.0, 1.5)
    moves = data.get('moves_2026', [])
    for i, line in enumerate(moves):
        c = GREEN if line.startswith('+') else (RED if line.startswith('-') else WHITE)
        add_text(s, 0.8, 1.3 + i * 0.35, 11.5, 0.35, line.lstrip('+-'), 13, c)
    watermark(s); slide_num(s, pc); pc += 1
    
    # --- SLIDE: Catalyst Timeline ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, 0.8, 0.4, 8, 0.6, f"{ticker} 2026 催化劑時間表", 30, RED, bold=True)
    add_gold_line(s, 0.8, 1.0, 1.5)
    catalysts = data.get('catalysts', [])
    for i, (time, event, impact) in enumerate(catalysts):
        y = 1.4 + i * 0.55
        add_card(s, 0.8, y, 11.5, 0.48)
        add_text(s, 1.0, y + 0.05, 1.5, 0.35, time, 12, GOLD, bold=True)
        add_text(s, 2.8, y + 0.05, 5.5, 0.35, event, 13, WHITE)
        c = GREEN if '利好' in impact or '催化' in impact or '最大' in impact else (RED if '風險' in impact else GRAY)
        add_text(s, 9.0, y + 0.05, 3, 0.35, impact, 11, c, align=PP_ALIGN.RIGHT)
    watermark(s); slide_num(s, pc); pc += 1
    
    # --- SLIDE: Special Focus (2 slides) ---
    focus_pages = data.get('focus', [])
    for fp in focus_pages:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s)
        add_text(s, 0.8, 0.4, 10, 0.6, fp['title'], 30, RED if 'risk' in fp.get('type','') else GOLD, bold=True)
        add_gold_line(s, 0.8, 1.0, 1.5)
        for i, line in enumerate(fp['lines']):
            c = GOLD if line.startswith('★') else (GREEN if line.startswith('+') else (RED if line.startswith('!') or line.startswith('-') else WHITE))
            txt = line.lstrip('★!+-')
            fs = 15 if line.startswith('##') else 13
            if line.startswith('##'):
                txt = txt.lstrip('# ')
                c = GOLD
                fs = 16
            add_text(s, 0.8, 1.3 + i * 0.35, 11.5, 0.35, txt, fs, c, bold=line.startswith('##') or line.startswith('★'))
        watermark(s); slide_num(s, pc); pc += 1
    
    # --- SLIDE: 競爭對手比較 (NEW) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, 0.8, 0.4, 8, 0.6, f"{ticker} 競爭對手比較", 30, BLUE, bold=True)
    add_gold_line(s, 0.8, 1.0, 1.5)
    
    competitors = data.get('competitors', [])
    for i, comp in enumerate(competitors):
        y = 1.4 + i * 1.15
        # Company name card
        add_card(s, 0.8, y, 11.5, 1.08)
        add_text(s, 1.0, y + 0.05, 10, 0.35, comp['name'], 14, GOLD, bold=True)
        add_text(s, 1.0, y + 0.42, 10.5, 0.3, comp['strategy_2026'], 12, WHITE)
        add_text(s, 1.0, y + 0.75, 10.5, 0.28, comp['advantage'], 11, GREEN if comp.get('winner') else GRAY)
    
    watermark(s); slide_num(s, pc); pc += 1
    
    # --- SLIDE: Verdict 投資結論 (CHANGED) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    
    # Decorative gold accent lines (thin vertical lines)
    for line_x in [0.5, 12.7]:
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(line_x), Inches(1.4), Pt(2), Inches(4.5))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.fill.transparency = 0.5
        line.line.fill.background()
    
    add_text(s, 0.8, 0.4, 8, 0.6, f"{ticker} 投資結論", 30, RED, bold=True)
    add_gold_line(s, 0.8, 1.0, 1.5)
    verdict = data.get('verdict', {})
    
    # Left: 公司最大優勢
    add_text(s, 0.8, 1.4, 5.5, 0.4, "公司最大優勢", 16, GREEN, bold=True)
    advantages = verdict.get('advantages', [])
    for i, adv in enumerate(advantages):
        add_text(s, 0.8, 1.9 + i * 0.35, 5.5, 0.35, f"▸ {adv}", 12, RGBColor(0xCC, 0xCC, 0xCC))
    
    # Right: 突破現況的最大難處
    add_text(s, 6.8, 1.4, 5.5, 0.4, "突破現況的最大難處", 16, RED, bold=True)
    challenges = verdict.get('challenges', [])
    for i, ch in enumerate(challenges):
        add_text(s, 6.8, 1.9 + i * 0.35, 5.5, 0.35, f"▸ {ch}", 12, RGBColor(0xCC, 0xCC, 0xCC))
    
    # Bottom: Key message
    add_card(s, 0.8, 5.5, 11.5, 1.0)
    add_text(s, 1.0, 5.6, 11, 0.7, verdict.get('key_message', ''), 15, GOLD, align=PP_ALIGN.CENTER)
    
    watermark(s); slide_num(s, pc); pc += 1
    
    page_counter[0] = pc


# ============================================================
# GOOGL DATA
# ============================================================
make_company_section('GOOGL', 'Alphabet Inc.', '成長瓶頸分析 + 2026突破可能性', '$303.33 · 市值$3.67T', '-2.4', {
    'key_stats': [
        ('2025 營收', '$403B', GOLD),
        ('YoY 成長', '+15.1%', GREEN),
        ('淨利率', '32.8%', WHITE),
        ('2025 CapEx', '$91B', RED),
    ],
    'overview': [
        '▸ 全球最大搜尋引擎 + 數位廣告平台，市佔率 90%+',
        '▸ CEO: Sundar Pichai · 員工: 190,820人',
        '▸ 三大業務板塊:',
        '    Google Services (90% 營收) — 搜尋廣告、YouTube、Google Play、Pixel',
        '    Google Cloud (10% 營收) — AI解決方案、Workspace、雲端基礎設施',
        '    Other Bets (<1%) — Waymo 自駕、Verily 健康科技',
        '',
        '▸ 2025年亮點:',
        '    營收首次突破 $4,000億，Q4 單季 $1,139億 創歷史新高',
        '    淨利 $1,322億 (+32% YoY)，淨利率提升至 32.8%',
        '    Q4 YoY 成長加速至 18%（Q1 僅 12%），趨勢向上',
    ],
    'revenue': [
        ('Search 廣告', 52, '~$210B (52%)', GOLD),
        ('YouTube 廣告', 9.5, '~$38B (9.5%)', RED),
        ('Google Cloud', 10, '~$40B (10%)', BLUE),
        ('Network/訂閱/其他', 7, '~$28B', GREEN),
        ('Other Bets', 1.5, '~$2B (<1%)', DIM),
    ],
    'capex_scenarios': [
        '★ 2025年 CapEx $910億 = 等於蓋 30 座台積電南科廠',
        '',
        '錢花在哪裡？（白話文解釋）',
        '▸ AI數據中心 (60-70%): 訓練Gemini大模型',
        '  → 想像成蓋一座城市大小的超級電腦機房',
        '  → 裡面塞滿 Google TPU 晶片 + NVIDIA GPU',
        '',
        '▸ 海底光纜 (5-10%): 全球網路骨幹',
        '  → Google 自己鋪設橫跨太平洋、大西洋的網路線',
        '  → 確保 YouTube、Gmail 全球秒開',
        '',
        '▸ 自駕車車隊 (Waymo): 造火箭就是燒錢',
        '  → 類似 Tesla 早期投資，但 Google 更激進',
        '  → 舊金山街頭跑的無人計程車',
        '',
        '!CapEx 成長 +74% 遠超營收成長 +15%，市場質疑 ROI',
        '!2026年 CapEx 預估: $750-850億（微降但仍極高）',
        '',
        '★ 核心問題: AI 軍備競賽的資本效率是否合理？',
        '  Google 的回答: "AI 投資不足的風險遠大於過度投資"',
    ],
    'moves_2026': [
        '+ 2月: Q4 財報超預期，營收 $1,139億 (+18% YoY)',
        '+ Gemini 2.0 持續迭代，在企業市場與 GPT-4o 競爭',
        '+ Google Cloud 2025年營收 ~$400億，營業利潤率改善至 8-10%',
        '',
        '- DOJ 反壟斷案: 2024年判決搜尋壟斷成立，2026年可能執行補救措施',
        '- AI Overviews 對廣告收入的衝擊: 每次 AI 搜尋廣告收入下降 15-20%',
        '- 第三方 Cookie 逐步淘汰，精準廣告能力下降',
        '',
        '+ Waymo 自駕: 2026年擴展至 LA、Austin 等新城市',
        '+ YouTube Shorts 日均觀看超 700億分鐘，廣告營利化加速',
        '+ Pixel 10 + Tensor G5 自研晶片，硬體生態持續建設',
    ],
    'catalysts': [
        ('3月', 'Google Cloud Next 大會', 'AI 新產品發布'),
        ('4月', 'Q1 2026 財報', 'AI 營收貢獻驗證'),
        ('5月', 'Google I/O — Gemini 3.0 發布', '年度最大催化劑'),
        ('Q2-Q3', 'AI Overviews 廣告上線', '搜尋營收突破口'),
        ('7月', 'Q2 2026 財報', 'Cloud 利潤率驗證'),
        ('下半年', 'DOJ 反壟斷補救措施', '潛在風險'),
        ('Q4', 'Waymo IPO/融資消息', '釋放隱藏價值'),
    ],
    'focus': [
        {
            'title': 'GOOGL 成長不足的真相',
            'type': 'analysis',
            'lines': [
                '## 問題一: 廣告業務增長見頂',
                '▸ 搜尋市佔率已達 90%+，缺乏增量市場',
                '▸ AI Overviews 直接回答問題，減少用戶點擊廣告',
                '▸ 測試數據: AI 搜尋頁面廣告收入下降 15-20%',
                '▸ TikTok、Amazon 搜尋、ChatGPT 持續分食市場',
                '',
                '## 問題二: Google Cloud 追趕困難',
                '▸ AWS +17%、Azure +32%、Google Cloud +21%',
                '▸ 企業客戶黏性不如 AWS/Azure 的 Office 365 生態系綁定',
                '▸ Cloud 營業利潤率僅 8-10%（AWS 30%+）',
                '',
                '## 問題三: CapEx 暴增但 ROI 存疑',
                '▸ 2023 CapEx $32B → 2025 $91B (+184%)，營收僅 +15%',
                '▸ CapEx/營收從 10.5% 暴增至 22.7%',
                '★ 核心矛盾: 搜尋是印鈔機，但 AI 正在瓦解它的商業模式',
            ]
        },
        {
            'title': 'GOOGL 2026 突破點在哪？',
            'type': 'analysis',
            'lines': [
                '## 突破點 1: AI Overviews 廣告化',
                '★ 在 AI 搜尋摘要中整合廣告，預估新增 $50-100億營收',
                '▸ 時間點: 2026 Q2-Q3 開始貢獻',
                '',
                '## 突破點 2: Google Cloud 利潤改善',
                '★ 2026 Cloud 目標: 營收 $500億+、營業利潤率 12-15%',
                '▸ AI 工作負載毛利率 50%+，拉高整體利潤',
                '',
                '## 突破點 3: YouTube Shorts 全面營利化',
                '★ Shorts 廣告佔 YouTube 總廣告 15-20%，新增 $30-50億',
                '',
                '## 突破點 4: Waymo 獨立估值',
                '★ 估值 $1,000億+，可能 2026 宣布 IPO',
                '',
                '+結論: 成長不足是暫時的。AI 是短期成本、長期武器。',
                '+2026 年是驗證 AI 投資回報的關鍵轉折年。',
            ]
        },
    ],
    'competitors': [
        {
            'name': 'GOOGL vs Microsoft (AI/Cloud)',
            'strategy_2026': 'Microsoft: Copilot 全線整合 Office 365，Azure AI 綁定企業客戶 | Google: Gemini 追趕，Cloud 價格戰搶市佔',
            'advantage': '✅ Microsoft 優勢: Office 生態系黏性 ❌ Google 劣勢: 企業市場起步晚',
            'winner': False
        },
        {
            'name': 'GOOGL vs Meta (廣告)',
            'strategy_2026': 'Meta: Reels 短影片廣告成熟，AI 個性化推薦 | Google: YouTube Shorts 追趕，搜尋廣告穩固',
            'advantage': '✅ Google 優勢: 搜尋意圖型廣告 ROI 高 ⚖️ Meta: 社群廣告年輕化',
            'winner': True
        },
        {
            'name': 'GOOGL vs OpenAI/ChatGPT',
            'strategy_2026': 'OpenAI: GPT-4o 生態系、API 收費 | Google: Gemini 免費 + 企業版，綁定 Workspace',
            'advantage': '⚖️ 技術平手，但 Google 分發管道更廣（Search、Android、Chrome）',
            'winner': True
        },
    ],
    'verdict': {
        'advantages': [
            '搜尋壟斷地位穩固（90%+ 市佔），護城河極深',
            '2025 淨利 $1,322億，財務狀況極度健康',
            'AI 技術領先（Gemini、TPU），不依賴 NVIDIA',
            'σ-2.4 超跌，P/E 28x 相對合理',
        ],
        'challenges': [
            'AI Overviews 可能瓦解搜尋廣告商業模式',
            'CapEx $910億 ROI 需要 2-3 年驗證',
            'DOJ 反壟斷判決可能限制搜尋業務結構',
            'Cloud 利潤率遠低於 AWS (8% vs 30%)',
        ],
        'key_message': '成長不足是暫時的 — AI 是短期成本、長期武器\n2026 Google I/O (5月) 是全年最大催化劑，σ-2.4 提供安全邊際'
    }
})

# ============================================================
# PLTR DATA
# ============================================================
make_company_section('PLTR', 'Palantir Technologies Inc.', '成長持續性 + 利潤率 + 2027方向', '$135.38 · 市值$3,093億', '-1.7', {
    'key_stats': [
        ('2025 營收', '$44.8B', GOLD),
        ('YoY 成長', '+56.2%', GREEN),
        ('淨利率', '36.3%', WHITE),
        ('毛利率', '82.4%', BLUE),
    ],
    'overview': [
        '▸ AI 驅動的資料整合與決策平台公司',
        '▸ CEO: Alex Karp · 員工: 4,001人 · 人均營收 $112萬（業界頂尖）',
        '▸ 核心產品:',
        '    Gotham — 情報/國防決策平台',
        '    Foundry — 企業級資料作業系統',
        '    AIP (AI Platform) — 2024年推出，爆發增長中',
        '',
        '▸ 2025年亮點:',
        '    營收 $44.75億 (+56.2%)，Q4 單季成長高達 70%',
        '    營業利潤率從 10.8% → 31.6%（+20.8個百分點）',
        '    淨利率從 16.1% → 36.3%，首次大規模獲利',
        '    自由現金流 $21億 (+84%)',
    ],
    'revenue': [
        ('美國政府', 40, '$16.2B (36%) +53%', RED),
        ('美國商業', 34, '$17.0B (38%) +74.7%', GOLD),
        ('國際政府', 14, '$7.8B (17.4%)', BLUE),
        ('國際商業', 8, '$3.7B (8.3%)', GREEN),
    ],
    'capex_scenarios': [
        '★ Palantir 是輕資產模式，錢花在人才上',
        '',
        '▸ 2025年 CapEx 僅 $1.5-2億（佔營收 <4%）',
        '  → 不像 Google/AMD 燒錢蓋廠房',
        '  → 主要支出: 雲端基礎設施、安全認證環境、辦公室',
        '',
        '真正的錢花在哪？',
        '▸ 軟體開發: 頂尖工程師年薪 $30-50萬美元',
        '▸ 安全認證: FedRAMP High、IL6 認證費用',
        '▸ 客戶成功團隊: AIP Bootcamp 全球巡迴',
        '',
        '營運計劃 2026:',
        '★ AIP 全面擴張 — 目標商業客戶 700+ (2025: ~500家)',
        '▸ AIP Bootcamp 轉換率提升，客戶平均合約金額從數十萬→數百萬',
        '▸ 國防合約鞏固: DoD、NATO、盟國長期戰略合作',
        '',
        '★ 營運槓桿驚人:',
        '  R&D/營收: 17.7% → 12.5% | S&M/營收: 31.0% → 23.6%',
        '  SBC/營收: 24.1% → 15.3%（股權稀釋壓力減緩）',
    ],
    'moves_2026': [
        '+ Q4 2025 營收成長 70% YoY，加速趨勢明顯',
        '+ AIP 客戶從「試點」→「企業級部署」，合約金額倍增',
        '+ 美國國防 AI 預算持續增長，Gotham 成為核心基礎設施',
        '+ 商業客戶增速超越政府客戶（+60% vs +53%）',
        '',
        '- 估值極高: P/E 214x、P/S 77x、EV/EBITDA 219x',
        '- Databricks、Microsoft Copilot 競爭加劇',
        '- 雙層股權結構: 創辦人持有 49.99% 投票權',
    ],
    'catalysts': [
        ('2月', 'Q4 2025 財報發布', '已超預期 (+70%)'),
        ('5月', 'Q1 2026 財報', '驗證成長持續性'),
        ('6月', 'AIPCon 2026 (預期)', '新產品發布 · 最大催化劑'),
        ('8月', 'Q2 2026 財報', '商業客戶增速追蹤'),
        ('Q3-Q4', '可能納入 S&P 500 指數', '被動資金大量流入'),
    ],
    'focus': [
        {
            'title': 'PLTR 成長率能否維持？',
            'type': 'analysis',
            'lines': [
                '## 2025: +56.2% → 2026E: +45-50%',
                '',
                '★ 支持高成長的因素:',
                '▸ AIP 動能: 商業客戶從 321→500+，合約金額倍增',
                '▸ 政府合約: 國防 AI 預算持續增長，續約率接近 100%',
                '▸ 規模效應: 營運利潤率 Q4 達 40.9%，固定成本攤提效應明顯',
                '',
                '!挑戰因素:',
                '▸ 基期效應: 2025 營收 $44.75B，維持 56% 需新增 $25B',
                '▸ 競爭: Databricks、Snowflake、Microsoft Azure AI',
                '',
                '+結論: 成長能維持。2026E +45-50%，仍遠超 SaaS 同業。',
                '+Q1-Q2 可能放緩至 40-45%，下半年 AIP 規模化後重返 50%+',
            ]
        },
        {
            'title': 'PLTR 利潤率 + 2027後方向',
            'type': 'analysis',
            'lines': [
                '## 利潤率還能提高嗎？',
                '★ 毛利率 82.4% → 接近天花板（軟體極限 85%）',
                '★ 營業利潤率 31.6% → 2026 目標 38-42%',
                '▸ 軟體邊際成本近零，每新增 $1 營收幾乎全是利潤',
                '▸ SBC 持續下降（24.1%→15.3%），利潤率改善空間仍大',
                '',
                '## 2027後的方向',
                '★ 短期 (2026-2027): 鞏固 AI 作業系統領導地位',
                '▸ 從「賣軟體」→「AI 作業系統平台」',
                '▸ 商業部門超越政府（2027 佔比 >55%）',
                '',
                '★ 中長期 (2028+): 平台生態系',
                '▸ 類似 Salesforce AppExchange 的第三方開發者平台',
                '▸ 數據市場 (Data Marketplace): 企業間數據交換',
                '▸ CEO Karp 願景: "數位神經系統，如 Windows 之於 PC"',
            ]
        },
    ],
    'competitors': [
        {
            'name': 'PLTR vs Databricks',
            'strategy_2026': 'Databricks: Lakehouse 平台，開源 Spark 生態 | PLTR: AIP 端到端整合，國防級安全',
            'advantage': '✅ PLTR 優勢: 政府客戶壁壘、作戰決策能力 ❌ Databricks: 數據科學家友好度更高',
            'winner': True
        },
        {
            'name': 'PLTR vs Snowflake + Microsoft Azure AI',
            'strategy_2026': 'Snowflake: 雲端數據倉儲易用 | Azure AI: Office 365 綁定 | PLTR: 垂直深度 + 安全認證',
            'advantage': '✅ PLTR 優勢: FedRAMP High + IL6 認證，競爭壁壘極高',
            'winner': True
        },
        {
            'name': 'PLTR vs C3.ai',
            'strategy_2026': 'C3.ai: 垂直產業 AI 應用豐富 | PLTR: 平台通用性 + 規模更大',
            'advantage': '✅ PLTR 優勢: 規模、客戶數、財務表現全面領先',
            'winner': True
        },
    ],
    'verdict': {
        'advantages': [
            'AIP 平台黏性極高，替換成本巨大',
            '營運槓桿驚人: 利潤率持續上升中',
            '國防 + 商業雙引擎，分散風險',
            'FedRAMP High + IL6 認證，競爭壁壘極高',
        ],
        'challenges': [
            '估值極高: P/E 214x，回調空間大',
            '雙層股權結構: 散戶話語權低',
            '客戶集中度偏高 (前10大客戶 30-40%)',
            'DCF 估值約 $90-110，當前溢價 23-50%',
        ],
        'key_message': '成長能維持 (45-50%)、利潤率能提高 (38-42%)\n最大風險是估值，但 AIP 平台是真正的護城河'
    }
})

# ============================================================
# RKLB DATA
# ============================================================
make_company_section('RKLB', 'Rocket Lab USA Inc.', '未來方向 + 政策面 + 成長率', '$74.42 · 市值$397.5億', '-0.6', {
    'key_stats': [
        ('2024 營收', '$436M', GOLD),
        ('YoY 成長', '+78%', GREEN),
        ('2024 淨虧損', '-$190M', RED),
        ('現金儲備', '$808M', BLUE),
    ],
    'overview': [
        '▸ 端到端太空公司，全球第二活躍火箭發射商（僅次於 SpaceX）',
        '▸ CEO: Peter Beck · 員工: ~2,100人',
        '▸ 核心業務:',
        '    Electron — 小型軌道火箭，80+ 次發射',
        '    Space Systems — Photon 衛星、航天零件、星座管理',
        '    Neutron — 8噸中型可重複使用火箭（開發中）',
        '',
        '▸ 2025年前三季:',
        '    營收 $422M (+38% YoY)，年化超過 $560M',
        '    毛利率從 29% → 37%，持續改善',
        '    Q3 虧損收窄至 $18M（Q1: $60M），盈利拐點接近',
    ],
    'revenue': [
        ('太空系統', 65, '~65% (衛星+零件)', GOLD),
        ('發射服務', 35, '~35% (Electron)', RED),
    ],
    'capex_scenarios': [
        '★ 造火箭就是燒錢（類似 SpaceX 早期）',
        '',
        '▸ 2024年 R&D: $174M（佔營收 40%）',
        '▸ Neutron 火箭累計投入 $150-200M',
        '',
        '2026 資本支出重點:',
        '▸ Neutron 火箭開發:',
        '  → 引擎測試: Archimedes 甲烷引擎全推力測試',
        '  → 發射台建設: 維吉尼亞州 LC-3 發射場',
        '  → 結構測試: 碳纖維機體壓力測試',
        '',
        '▸ Electron 產線優化:',
        '  → 提高發射頻率至 15-18 次/年',
        '  → 回收再利用技術測試（降低成本）',
        '',
        '▸ 衛星生產線擴建:',
        '  → Mandrake-2 後續訂單交付',
        '  → Photon 衛星平台標準化',
        '',
        '★ 現金 $808M，足以支撐至 2027 年中',
        '!若 Neutron 延遲或成本超支，可能需要再融資',
    ],
    'moves_2026': [
        '+ Mandrake-2 太空機動任務成功，展示軍事應用能力',
        '+ Electron 第 80 次發射達成，可靠性持續驗證',
        '+ 毛利率改善至 37%（2022年為 9%），軌跡向好',
        '+ SDA (太空發展局) 星座計劃 $200億+ 預算',
        '',
        '- 淨虧損持續（2024: -$190M），尚未盈利',
        '- Neutron 已多次推遲，首飛時間不確定',
        '- 估值 EV/Sales 44x，需要高成長來支撐',
        '- SpaceX 壟斷 70%+ 發射市場',
    ],
    'catalysts': [
        ('Q1', '2025 Q4 財報', '年化營收 $600M+ 驗證'),
        ('Q2', 'SDA 星座計劃競標結果', '政府業務爆發點'),
        ('Q3', 'Neutron 引擎全量測試', '技術風險消除'),
        ('Q4', 'Neutron 首飛', '年度最大催化劑'),
        ('全年', 'Electron 15-18 次發射', '可靠性持續驗證'),
    ],
    'focus': [
        {
            'title': 'RKLB 未來方向與計劃',
            'type': 'analysis',
            'lines': [
                '## 短期 (2026): Neutron 首飛是一切',
                '★ Neutron 成功 → 解鎖 $200B+ 中型火箭市場',
                '▸ 8噸載荷，可重複使用，直接與 Falcon 9 競爭',
                '',
                '## 中期 (2027-2029): 雙引擎加速',
                '★ 發射服務: Electron + Neutron 並行，年營收 $500M+',
                '★ 太空系統: 衛星平台 + 零件，年營收 $800M+',
                '',
                '## 長期 (2030+): 太空基礎設施巨頭',
                '▸ 全球第二大太空公司（僅次於 SpaceX）的目標',
                '',
                '+成長率會不會重大改變？',
                '+2026E +60%，Neutron 成功後 2027+ 可能加速至 CAGR 40-45%',
                '!失敗風險: Neutron 若延遲至 2027，成長率可能降至 20-25%',
            ]
        },
        {
            'title': 'RKLB 政策面分析',
            'type': 'analysis',
            'lines': [
                '## 順風: 政府太空預算持續增長',
                '★ Space Force 預算 $30-33B/年，持續增長中',
                '★ SDA 星座計劃 $200億+（5年期），RKLB 是主要競標者',
                '▸ NASA 預算 $25-27B，穩定',
                '▸ 盟國 (日本、澳洲、英國) 增加太空國防支出',
                '',
                '## 政府業務佔比: 預估 40-50% 營收',
                '▸ 政策貢獻: 預計 2026-2030 提升成長率 10-15%',
                '',
                '★ 結論: 政策面整體強烈順風',
                '▸ 太空國防是兩黨共識，不受政黨輪替影響',
            ]
        },
    ],
    'competitors': [
        {
            'name': 'RKLB vs SpaceX',
            'strategy_2026': 'SpaceX: Falcon 9/Heavy 主導市場，Starship 開發 | RKLB: Neutron 填補中型火箭空缺',
            'advantage': '❌ SpaceX 絕對優勢: 成本、頻率、規模 ✅ RKLB 機會: 政府需要備選方案',
            'winner': False
        },
        {
            'name': 'RKLB vs Firefly + Relativity Space',
            'strategy_2026': 'Firefly: Alpha 火箭競爭小型市場 | Relativity: Terran R 比 Neutron 更大 | RKLB: 技術成熟度領先',
            'advantage': '✅ RKLB 優勢: 80+ 次發射記錄，競爭者仍在驗證階段',
            'winner': True
        },
        {
            'name': 'RKLB vs 中國火箭公司（長征、星際榮耀）',
            'strategy_2026': '中國: 成本更低，但受 ITAR 限制無法進入美國市場 | RKLB: 美國本土製造優勢',
            'advantage': '✅ RKLB 優勢: 國防合約、安全供應鏈，中國無法競爭',
            'winner': True
        },
    ],
    'verdict': {
        'advantages': [
            'Neutron 解鎖 $200B+ 市場機會',
            '太空國防預算持續增長（兩黨共識）',
            '垂直整合優勢，成本控制佳',
            '毛利率從 9% → 37%，盈利拐點接近',
        ],
        'challenges': [
            'Neutron 延遲風險（已推遲數次）',
            '估值偏高 EV/Sales 44x',
            '現金流為負，支撐至 2027 中',
            'SpaceX 壟斷壓力',
        ],
        'key_message': 'Neutron 首飛是 2026 最大催化劑\n成功 → 太空版 Tesla 故事成型 | 失敗 → 估值大幅修正'
    }
})

# ============================================================
# AMD DATA
# ============================================================
make_company_section('AMD', 'Advanced Micro Devices Inc.', '低利潤率突破路徑', '$200.12 · 市值$3,250億', '-2.2', {
    'key_stats': [
        ('2025 營收', '$34.6B', GOLD),
        ('YoY 成長', '+34.3%', GREEN),
        ('淨利率', '12.5%', YELLOW),
        ('自由現金流', '$6.7B', BLUE),
    ],
    'overview': [
        '▸ 全球第二大 x86 處理器廠商，NVIDIA 最大 AI GPU 挑戰者',
        '▸ CEO: Lisa Su · 員工: ~26,000人',
        '▸ 核心產品線:',
        '    EPYC — 伺服器 CPU（市佔 25-30%）',
        '    MI300X/MI350/MI400 — AI GPU',
        '    Ryzen — 消費端 CPU + AI PC',
        '',
        '▸ 2025年亮點:',
        '    營收 $346億 (+34.3%)，數據中心成為核心引擎',
        '    淨利 $43.35億 (+164%)，Q4 毛利率 54.3%（歷史新高）',
        '    自由現金流 $67億 (+180%)',
        '',
        '★ 核心矛盾: 淨利率 12.5% vs NVIDIA 55%，差距 4 倍',
    ],
    'revenue': [
        ('數據中心', 48, '$166B (48%) +32%', GOLD),
        ('客戶端', 30, '~$104B (~30%) +106%', BLUE),
        ('遊戲', 12, '~$42B (~12%)', RED),
        ('嵌入式', 10, '$34.5B (10%) -2.9%', GREEN),
    ],
    'capex_scenarios': [
        '★ AMD 是 Fabless 無廠，錢花在設計而非硬體',
        '',
        '▸ 2025年 R&D: ~$55-60億（佔營收 16-17%）',
        '▸ 2025年 CapEx: ~$7-8億（測試設備、辦公設施）',
        '',
        '錢花在哪裡？',
        '▸ AI晶片設計（對打 NVIDIA）:',
        '  → MI350/MI400 AI GPU — 挑戰 GB200',
        '  → EPYC Turin (Zen 5) — 192核伺服器 CPU',
        '',
        '▸ 軟體生態（ROCm 追趕 CUDA）:',
        '  → ROCm 6.0+ 開源軟體堆疊',
        '  → 企業客戶技術支援團隊',
        '',
        '▸ 台積電代工費:',
        '  → N5/N4P 晶圓 $16,000-18,000/片',
        '  → HBM3/HBM3E 記憶體採購（單顆 $200-300）',
        '',
        '★ 資本支出效率高（Fabless 模式），錢花在刀口上',
        '!最大支出風險: 台積電先進製程價格上漲',
    ],
    'moves_2026': [
        '+ Q4 2025 毛利率 54.3% 創歷史新高，突破 50% 里程碑',
        '+ MI300X AI GPU 獲微軟、Meta 等大客戶採用',
        '+ EPYC 伺服器 CPU 市佔持續爬升至 25-30%',
        '+ 自由現金流 $67億，財務體質大幅改善',
        '',
        '- NVIDIA 壟斷 AI GPU 80%+ 市場，AMD 僅 10-15%',
        '- ROCm 軟體生態仍遠落後 CUDA',
        '- HBM (高頻寬記憶體) 成本壓力: 佔 AI GPU 成本 40-50%',
        '- 遊戲部門萎縮中 (-2.9%)，低毛利拖累整體',
    ],
    'catalysts': [
        ('Q1-Q2', 'MI350 AI GPU 量產出貨', 'AI GPU 營收加速'),
        ('Q2', 'EPYC Turin (Zen 5) 全面鋪貨', 'CPU 市佔突破'),
        ('Q3-Q4', 'MI400 AI GPU 發布', '年度最大催化劑'),
        ('全年', 'AI PC 滲透率提升 (Ryzen AI)', '客戶端成長'),
        ('全年', 'ROCm 6.0+ 軟體更新', '生態追趕 CUDA'),
    ],
    'focus': [
        {
            'title': 'AMD 低利潤率的根本原因',
            'type': 'analysis',
            'lines': [
                '## AMD 淨利率 12.5% vs NVIDIA 55% — 為什麼？',
                '',
                '!原因 1: 產品組合拖累',
                '▸ 低毛利 Client CPU (40-45%) + Gaming (30-35%) 佔比過高',
                '▸ 高毛利 AI GPU 佔總營收僅 15-20%',
                '',
                '!原因 2: HBM 成本壓力',
                '▸ MI300X 使用 8 顆 HBM3，單顆 $200-300',
                '▸ HBM 佔 AI GPU 成本 40-50%，壓縮毛利率',
                '',
                '!原因 3: 規模效應不足',
                '▸ 訂單量遠小於 NVIDIA，台積電/HBM 議價能力弱',
                '',
                '!原因 4: 競爭定價',
                '▸ EPYC 為搶 Intel 市佔，價格折讓 10-15%',
                '▸ MI300X 走性價比路線，售價低 NVIDIA 20-30%',
            ]
        },
        {
            'title': 'AMD 利潤率突破路徑',
            'type': 'analysis',
            'lines': [
                '## 三種情境預測',
                '',
                '★ 樂觀 (2027 毛利率 60%+):',
                '▸ AI GPU 佔營收 40%+、毛利率 65-70%',
                '▸ EPYC 市佔 35%+，規模效應顯現',
                '▸ HBM 成本下降 15-20%',
                '▸ → 淨利率 25%',
                '',
                '+中性 (2027 毛利率 55%):',
                '▸ AI GPU 佔營收 30%、毛利率 65%',
                '▸ EPYC 市佔 30%',
                '▸ → 淨利率 18-20%',
                '',
                '!悲觀 (毛利率停滯 50-52%):',
                '▸ NVIDIA 壟斷持續、Intel 反撲',
                '▸ HBM 供應短缺漲價',
                '',
                '★ 關鍵拐點: 2026 Q4 MI400 發布 → 毛利率推至 54-55%',
                '★ 2027 全年毛利率 57-60%，淨利率 20%',
            ]
        },
    ],
    'competitors': [
        {
            'name': 'AMD vs NVIDIA (AI GPU)',
            'strategy_2026': 'NVIDIA: GB200/Blackwell Ultra 生態系壁壘 | AMD: MI400 性價比路線，ROCm 追趕',
            'advantage': '❌ NVIDIA 壓倒性優勢: CUDA 生態、80%+ 市佔 ✅ AMD 機會: 雲端客戶需要第二供應商',
            'winner': False
        },
        {
            'name': 'AMD vs Intel (伺服器 CPU)',
            'strategy_2026': 'Intel: Granite Rapids/Sierra Forest 反撲 | AMD: EPYC Turin (Zen 5) 192核優勢',
            'advantage': '✅ AMD 優勢: 市佔率持續擴張 25%→35%，Intel 製程落後',
            'winner': True
        },
        {
            'name': 'AMD vs Qualcomm/Apple (AI PC)',
            'strategy_2026': 'Qualcomm: ARM 架構 AI PC | Apple: M4 自研晶片 | AMD: Ryzen AI 300 x86 相容性',
            'advantage': '⚖️ AMD 優勢: Windows 生態相容 ❌ 劣勢: 能效比不如 ARM',
            'winner': True
        },
    ],
    'verdict': {
        'advantages': [
            'AI GPU 唯一能挑戰 NVIDIA 的對手',
            'EPYC 持續搶 Intel 市佔 (25%→35%)',
            'Q4 毛利率 54.3% 已開始突破',
            'σ-2.2 超跌，PEG 0.49（便宜）',
        ],
        'challenges': [
            'NVIDIA 壟斷 AI GPU 80%+，CUDA 生態壁壘高',
            'HBM 成本壓力短期難解',
            'ROCm 軟體生態仍需 2-3 年追趕',
            '遊戲部門萎縮，低毛利拖累',
        ],
        'key_message': '利潤率拐點就在 2026 Q4 — MI400 發布是關鍵\nσ-2.2 + PEG 0.49 = 被低估的 NVIDIA 挑戰者'
    }
})

# ============================================================
# ARES DATA
# ============================================================
make_company_section('ARES', 'Ares Management Corporation', '政策影響 + 寬鬆/緊縮情境', '$133.97 · 市值$490億', '-1.6', {
    'key_stats': [
        ('2025 營收', '~$70B', GOLD),
        ('YoY 成長', '+66.6%', GREEN),
        ('AUM', '~$5,000B', BLUE),
        ('殖利率', '2.77%', WHITE),
    ],
    'overview': [
        '▸ 全球頂尖另類資產管理公司，私募信貸領域全球 #1',
        '▸ CEO: Michael Arougheti · 員工: ~3,000人',
        '▸ 業務板塊:',
        '    信貸/直接貸款 (~60%) — 核心引擎',
        '    私募股權 (~15%) — 中型企業收購',
        '    房地產 (~15%) — 商業地產債權與股權',
        '',
        '▸ 2025年亮點:',
        '    營收暴增 +66.6%，受私募信貸爆發驅動',
        '    AUM 估計 $4,500-5,000億，持續擴張',
        '    輕資產模式: CapEx < 3% 營收',
        '    管理費佔收入 70-75%（穩定現金流）',
    ],
    'revenue': [
        ('信貸/直接貸款', 60, '~60% (核心引擎)', GOLD),
        ('私募股權', 15, '~15%', BLUE),
        ('房地產', 15, '~15%', RED),
        ('基礎設施', 10, '~10%', GREEN),
    ],
    'capex_scenarios': [
        '★ Ares 幾乎不花錢在硬體（輕資產模式）',
        '',
        '▸ CapEx 佔營收 < 3%',
        '  → 主要支出: IT 系統、辦公設施、合規基礎設施',
        '  → 不像製造業或科技公司需要蓋廠房',
        '',
        '真正的錢花在哪？（替客戶投資）',
        '★ 2025年新增部署: ~$500-700億（私募信貸為主）',
        '',
        '▸ 直接貸款（替代銀行放貸）:',
        '  → 中端市場企業融資，利率 SOFR + 4-6%',
        '  → 客戶: EBITDA $1,000萬-$2.5億的企業',
        '',
        '▸ 基礎設施投資:',
        '  → 數據中心、再生能源項目、5G基地台',
        '',
        '▸ 房地產機會型投資:',
        '  → 不良資產收購（商業地產市場承壓）',
        '',
        '2026 營運計劃:',
        '▸ 新基金募資: 大型直接貸款基金 + 基礎設施基金',
        '▸ 零售化: Evergreen 基金吸引高淨值個人',
    ],
    'moves_2026': [
        '+ 私募信貸結構性增長: $1.5T → $3.5T by 2030',
        '+ 銀行退出中端市場 = Ares 最大紅利',
        '+ 管理費收入穩定增長，不受市場波動影響',
        '+ 機構投資者持續增配另類資產 (5-10%→15-20%)',
        '',
        '- 2025 基期 +66.6% 太高，2026 成長必然放緩至 20-30%',
        '- 派息比率 184.7%（超過淨利），可持續性存疑',
        '- SEC 對私募基金費用透明度監管加強',
        '- ROE 下降至 6.1%（2023: 25.1%），需關注盈利能力',
    ],
    'catalysts': [
        ('Q1', 'Q4 2025 財報 + 2026 指引', 'AUM 增長驗證'),
        ('2-3月', '大型直接貸款基金完成募資', '管理費收入提升'),
        ('Q2', 'Q1 季報', '私募信貸部署速度'),
        ('5-6月', '潛在併購公告', '業務規模擴張'),
        ('Q3', '聯準會利率政策轉折點', '影響全局 · 最大催化劑'),
    ],
    'focus': [
        {
            'title': 'ARES 成長能否持續擴張？',
            'type': 'analysis',
            'lines': [
                '## 結構性增長 vs 基期效應',
                '',
                '★ 支持持續增長:',
                '▸ 私募信貸市場 $1.5T → $3.5T by 2030（CAGR ~15%）',
                '▸ 銀行退出中端市場，監管收緊 → 企業找 Ares',
                '▸ 機構投資者增配另類資產趨勢不變',
                '▸ Evergreen 基金開拓零售渠道',
                '',
                '!增長放緩風險:',
                '▸ 2025 +66.6% 基期太高（含一次性併購/會計調整）',
                '▸ 2026 預估回落至 +15-25%',
                '▸ 私募信貸市場競爭加劇，管理費率面臨下行壓力',
                '',
                '+結論: 長期 CAGR 20%+ 可期，但短期增速必然回落',
                '+Ares 在私募信貸的領導地位是最大結構性優勢',
            ]
        },
        {
            'title': 'ARES 貨幣政策情境分析',
            'type': 'analysis',
            'lines': [
                '## 為何高利率反而有利？',
                '▸ 直接貸款利率隨 SOFR 浮動 → 高利率 = 高利息收入',
                '▸ 銀行收緊放貸 → 更多企業找 Ares',
                '',
                '## 三種情境:',
                '',
                '+情境一: 降息至 3.5-4%（軟著陸）— 最佳',
                '▸ 營收 $80.5B (+15%) · EPS $3.25 (+30%)',
                '▸ PE/RE 估值重估 → 目標 $187-200 (+40-50%)',
                '',
                '★ 情境二: 維持 5%+ 高利率 — 中性',
                '▸ 營收 $73.5B (+5%) · EPS $2.13 (-15%)',
                '▸ 直接貸款仍強 → 目標 $107-120 (-10-20%)',
                '',
                '!情境三: 硬著陸 + 降息 — 最差',
                '▸ 營收 $66.5B (-5%) · EPS $1.38 (-45%)',
                '▸ 違約率飆升 → 目標 $73-87 (-35-45%)',
                '',
                '★ 最佳情境: 軟著陸 + 緩降息 → $200+',
            ]
        },
    ],
    'competitors': [
        {
            'name': 'ARES vs Blackstone (全方位競爭)',
            'strategy_2026': 'Blackstone: 規模最大 ($1萬億+ AUM)，多元化 | Ares: 私募信貸專注，中端市場深耕',
            'advantage': '❌ Blackstone 規模優勢 ✅ Ares 私募信貸市佔率更高',
            'winner': False
        },
        {
            'name': 'ARES vs Apollo (信貸平台競爭)',
            'strategy_2026': 'Apollo: 信貸 + 保險整合 (Athene) | Ares: 純資產管理模式',
            'advantage': '⚖️ Apollo 保險資金穩定 vs Ares 靈活度高',
            'winner': True
        },
        {
            'name': 'ARES vs KKR (私募股權傳統強項)',
            'strategy_2026': 'KKR: PE 傳統巨頭，信貸後發 | Ares: 信貸領先，PE 相對弱',
            'advantage': '✅ Ares 信貸優勢明顯，但 PE 業務落後',
            'winner': True
        },
    ],
    'verdict': {
        'advantages': [
            '私募信貸全球 #1，結構性增長趨勢',
            '輕資產模式，管理費收入穩定',
            '軟著陸情境下上漲空間 40-50%',
            '殖利率 2.77%，提供下檔保護',
        ],
        'challenges': [
            '硬著陸風險: 違約率飆升 → -35-45%',
            '2025 基期太高，2026 增速必然放緩',
            'ROE 下降至 6.1%，盈利效率存疑',
            '派息比率 184.7% > 100%，不可持續',
        ],
        'key_message': '最大的 beta 來自聯準會 — 降息確認是最強催化劑\n私募信貸結構性增長是長期α，但需承受利率政策的短期波動'
    }
})

# ============================================================
# SUMMARY SLIDES
# ============================================================

# Slide: 5-stock comparison
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.4, 8, 0.6, "五檔反市場精選 — 綜合比較", 30, GOLD, bold=True)
add_gold_line(s, 0.8, 1.0, 2)

# Table header
headers = ['股票', 'σ', 'R40', '2026E成長', '核心看點', '最大風險', '催化時間']
for i, h in enumerate(headers):
    x = 0.5 + i * 1.8
    w = 1.6 if i < 3 else 2.0
    if i >= 4: x = 0.5 + 3*1.8 + (i-3)*2.5; w = 2.3
    add_text(s, x, 1.3, w, 0.3, h, 9, GOLD, bold=True)

rows = [
    ('GOOGL', '-2.4', '47', '+14%', 'AI Overviews 廣告化', 'CapEx ROI · 反壟斷', '5月 I/O'),
    ('PLTR', '-1.7', '92.5', '+45-50%', 'AIP 平台黏性', '估值 P/E 214x', '持續加速'),
    ('RKLB', '-0.6', '34.7', '+60%', 'Neutron 首飛', '虧損 · 延遲', 'Q4 首飛'),
    ('AMD', '-2.2', '46.9', '+30-36%', 'MI400 + 利潤拐點', 'NVIDIA 壟斷', 'Q3-Q4 MI400'),
    ('ARES', '-1.6', '74.7', '+15-25%', '私募信貸結構增長', '硬著陸風險', '降息確認'),
]
for ri, row in enumerate(rows):
    y = 1.7 + ri * 0.65
    if ri % 2 == 0:
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(y - 0.05), Inches(12.5), Inches(0.55))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x14,0x14,0x14); shape.line.fill.background()
    for ci, val in enumerate(row):
        x = 0.5 + ci * 1.8
        w = 1.6 if ci < 3 else 2.0
        if ci >= 4: x = 0.5 + 3*1.8 + (ci-3)*2.5; w = 2.3
        c = GOLD if ci == 0 else (RED if '-' in val and ci == 1 else (GREEN if '+' in val else WHITE))
        add_text(s, x, y, w, 0.4, val, 11, c, bold=(ci == 0))

# Bottom cards
add_card(s, 0.5, 5.2, 5.8, 0.9)
add_text(s, 0.7, 5.25, 5.4, 0.3, '最深超跌（技術面機會）', 11, RED, bold=True)
add_text(s, 0.7, 5.6, 5.4, 0.35, 'GOOGL (-2.4σ) > AMD (-2.2σ) > PLTR (-1.7σ)', 15, GOLD, bold=True)

add_card(s, 6.8, 5.2, 5.8, 0.9)
add_text(s, 7.0, 5.25, 5.4, 0.3, '最強基本面（R40）', 11, GREEN, bold=True)
add_text(s, 7.0, 5.6, 5.4, 0.35, 'PLTR (92.5) > ARES (74.7) > GOOGL (47)', 15, GOLD, bold=True)

watermark(s)
pc = page_counter[0]
slide_num(s, pc); pc += 1

# Final slide: Disclaimer
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1.5, 2.0, 10, 1.0, "JG的反市場報告書", 48, GOLD, bold=True, align=PP_ALIGN.CENTER)
add_gold_line(s, 5.5, 3.2, 2.3)
add_text(s, 1.5, 3.8, 10, 0.5, "深度研究系列 · 2026年2月19日", 16, GRAY, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 4.8, 10, 0.8, "本報告僅供參考，不構成任何投資建議。\n投資有風險，入市需謹慎。過往表現不代表未來收益。", 12, DIM, align=PP_ALIGN.CENTER)
watermark(s)
slide_num(s, pc)

# Save
output_path = os.path.join(os.path.dirname(__file__), 'JG反市場報告書_深度研究.pptx')
prs.save(output_path)
print(f'✅ PPT saved to: {output_path}')
print(f'📊 Total slides: {len(prs.slides)}')
